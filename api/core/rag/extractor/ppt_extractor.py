"""Abstract interface for document loader implementations."""

from core.rag.extractor.extractor_base import BaseExtractor
from core.rag.models.document import Document
from pptx import Presentation
import re
import uuid
import os
import shutil

class PptExtractor(BaseExtractor):
    """Load Excel files.


    Args:
        file_path: Path to the file to load.
    """

    def __init__(
            self,
            file_path: str,
    ):
        """Initialize with file path."""
        self._file_path = file_path

    def extract(self) -> list[Document]:
        if self._file_path.endswith('.ppt'):
            return self._extract4ppt(self._file_path)
        elif self._file_path.endswith('.pptx'):
            return self._extract4pptx(self._file_path)

    def _extract4ppt(self,fp) -> list[Document]:
        # first install :  sudo apt-get update
        # sudo apt-get install libreoffice
        # libreoffice --headless --convert-to pptx  haha.ppt
        tmp_folder = str(uuid.uuid1())
        cmd = f'libreoffice --headless --convert-to pptx {fp} --outdir {tmp_folder}'
        os.system(cmd)
        fs = os.listdir(tmp_folder)
        for f in fs:
            if f.endswith("pptx"):
                ds = self._extract4pptx(os.path.join(tmp_folder,f))
                shutil.rmtree(tmp_folder)
                return ds
        return []



    def _extract4pptx(self,fp) -> list[Document]:
        data = []
        # 打开一个PPT文件
        prs = Presentation(fp)
        for slide_num, slide in enumerate(prs.slides, start=1):
            # 遍历幻灯片中的每一个形状
            item = ""
            for shape_num, shape in enumerate(slide.shapes, start=1):
                # 如果形状有文本框
                if hasattr(shape, "text") and shape.text.strip():
                    item = item + shape.text.strip() + '\n'

            if len(item) > 1:
                item = re.sub(r'[\n]+', '\n', item)  # 页内尽量不切开
                document = Document(page_content=item, metadata={'source': self._file_path})
                data.append(document)

        return data


